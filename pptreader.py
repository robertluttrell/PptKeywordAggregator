import os
from pptx import Presentation
from keywordfilepresence import KeywordFilePresence


class PptReader:

    def __init__(self, ppt_path_list):
        self.ppt_path_list = ppt_path_list
        self.word_dict = {}

    def keyword_file_presence_exists(self, keyword, file_path):
        """
        Returns true if the KeywordFilePresence exists, indicating there's already an entry for this keyword and slide
        :param keyword: String representing the keyword
        :param file_path: String representing the path to the PowerPoint file the keyword was found in
        :return: True if exists, else false
        """
        keyword_kfp_list = self.word_dict[keyword]
        for kfp in keyword_kfp_list:
            if kfp.file_path == file_path:
                return True

        return False

    def process_keyword(self, keyword, index, file_path):
        """
        Adds the KeywordFilePresence to self.word_dict
        :param keyword: String representing the keyword
        :param index: Index of slide in powerpoint presentation starting from 1
        :param file_path: String representing path to the PowerPoint presentation containing the slide
        :return:
        """
        if keyword not in self.word_dict:
            self.word_dict[keyword] = []

        word_dict_entry = self.word_dict[keyword]

        if not self.keyword_file_presence_exists(keyword, file_path):
            word_dict_entry.append(KeywordFilePresence(file_path))

        word_dict_entry[-1].add_index(index)

    @staticmethod
    def remove_empty_strings(string_list):
        return list(filter(None, string_list))

    def process_slide(self, slide, index, file_path):
        """
        Reads a slide and stores keyword data in self.word_dict if a populated notes slide exists
        :param slide: Slide object to be processed
        :param index: Index of slide in powerpoint presentation starting from 1
        :param file_path: String representing path to the PowerPoint presentation containing the slide
        """
        keyword_text = None

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.startswith("Keyword"):
                keyword_text = shape.text.split(':', 1)[-1]
                break

        if keyword_text is None:
            return

        slide_keyword_list = self.remove_empty_strings(keyword_text.split(','))

        for i in range(len(slide_keyword_list)):
            slide_keyword_list[i] = slide_keyword_list[i].strip(' ').replace('\n', '')

        for keyword in slide_keyword_list:
            self.process_keyword(keyword, index, file_path)

    def ppt_file_to_dict(self, file_path):
        """
        Reads the file and stores keyword data in self.word_dict
        :param file_path: String path to the file to be read
        """
        try:
            file = open(file_path, "rb")

        except IOError as e:
            print(e)
            return

        pres = Presentation(file)
        file.close()

        for i in range(len(pres.slides)):
            self.process_slide(pres.slides[i], i + 1, file_path)

    def ppt_files_to_dict(self):
        """
        For each file specified in ppt_path_list, reads the file and stores keyword data in self.word_dict
        """
        if len(self.ppt_path_list) == 0:
            return

        for file_path in self.ppt_path_list:
            self.ppt_file_to_dict(file_path)
